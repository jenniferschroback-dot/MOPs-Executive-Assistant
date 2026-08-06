"""
Reusable python-pptx builders implementing the Acquia brand system and layout
templates documented in reference.md. Copy/adapt this module per-deck rather
than importing it as a shared library across unrelated runs (keeps each deck
self-contained and easy to hand-edit).

Usage pattern:
    from pptx_helpers import *
    prs = new_presentation()
    add_title_slide(prs, "Pardot Email Performance", "Last 2 Weeks", "July 3-17, 2026 | MOPS")
    add_overview_cards_slide(prs, "Core Metrics", "Email Performance", "Intro sentence.",
        cards=[(ICONS["growth-chart"], "Open Rate", "10.66%, +2.0 pts vs prior"), ...])
    ...
    prs.save("outputs/decks/My_Deck.pptx")   # outputs/ is one subfolder per subject —
                                             # see .claude/rules/output-files.md
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
ICONS_DIR = os.path.join(ASSETS, "icons")
ICONS = {name[:-4]: os.path.join(ICONS_DIR, name) for name in os.listdir(ICONS_DIR)} if os.path.isdir(ICONS_DIR) else {}
LOGO_BLUE = os.path.join(ASSETS, "logo-blue.png")
LOGO_NAVY = os.path.join(ASSETS, "logo-navy.png")
DROPLET = os.path.join(ASSETS, "droplet-blue.png")
PATTERN_SUBTLE = os.path.join(ASSETS, "pattern-subtle.png")

# ---- Brand colors (see reference.md for provenance) ----
NAVY = RGBColor(0x23, 0x2C, 0x61)
ACQUIA_BLUE = RGBColor(0x03, 0x6B, 0xB5)
SKY_BLUE = RGBColor(0x26, 0xA3, 0xDD)
BODY_GRAY = RGBColor(0x3D, 0x4F, 0x5C)
CARD_FILL = RGBColor(0xF2, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# Secondary/reserve accents (available, unused by default)
TEAL = RGBColor(0x66, 0xC8, 0xCA)
PINK = RGBColor(0xE1, 0x12, 0x6E)
ORANGE = RGBColor(0xF4, 0x7A, 0x20)
YELLOW = RGBColor(0xFE, 0xC2, 0x31)

# ---- Fonts: Montserrat by default (free, Google-Slides-native match for
# Proxima Nova). Switch to Proxima Nova only if generating for a machine
# that has it installed/licensed. ----
FONT_HEADING = "Montserrat ExtraBold"
FONT_BODY = "Montserrat"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SCALE = 13.333 / 10  # source deck was authored on a 10x5.625in canvas


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _run(p, text, size, color, bold=False, italic=False, font=None):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font or FONT_BODY
    return r


def add_logo(slide, variant="blue", top=Inches(0.5), left=Inches(0.5), width=Inches(1.5)):
    path = LOGO_BLUE if variant == "blue" else LOGO_NAVY
    slide.shapes.add_picture(path, left, top, width=width)


def add_footnote(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    _run(p, text, 8.5, BODY_GRAY)


def _eyebrow_title_intro(slide, eyebrow, title, intro):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.5), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    _run(p, eyebrow.upper(), 12, NAVY, bold=True, font=FONT_BODY)

    tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.75), Inches(12.2), Inches(0.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    _run(p2, title, 28, NAVY, bold=True, font=FONT_HEADING)

    if intro:
        tb3 = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(11.8), Inches(0.7))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        _run(p3, intro, 14, BODY_GRAY)


# ---------------------------------------------------------------------------
# Template 1 / 7: Title / section-break / closing slide (navy + droplets)
# ---------------------------------------------------------------------------
def add_title_slide(prs, title, subtitle="", footer="", eyebrow_logo="blue"):
    slide = _blank(prs)
    _set_bg(slide, NAVY)
    if os.path.exists(PATTERN_SUBTLE):
        slide.shapes.add_picture(PATTERN_SUBTLE, 0, 0, width=SLIDE_W, height=SLIDE_H)
    if os.path.exists(DROPLET):
        slide.shapes.add_picture(DROPLET, Inches(7.75 * SCALE), Inches(-0.55 * SCALE),
                                  width=Inches(2.55 * SCALE), height=Inches(3.31 * SCALE))
        slide.shapes.add_picture(DROPLET, Inches(7.75 * SCALE), Inches(2.76 * SCALE),
                                  width=Inches(2.55 * SCALE), height=Inches(3.31 * SCALE))
    add_logo(slide, variant=eyebrow_logo, top=Inches(0.6), left=Inches(0.6), width=Inches(1.8))

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(9.5), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _run(p, title, 40, WHITE, bold=True, font=FONT_HEADING)
    if subtitle:
        p2 = tf.add_paragraph()
        _run(p2, subtitle, 20, SKY_BLUE)
    if footer:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(18)
        _run(p3, footer, 13, RGBColor(0xC7, 0xD6, 0xE3))
    return slide


add_closing_slide = add_title_slide  # same template


# ---------------------------------------------------------------------------
# Template 2: Overview + N-card row (2-4 cards)
# ---------------------------------------------------------------------------
def add_overview_cards_slide(prs, eyebrow, title, intro, cards, logo_variant="blue"):
    """cards: list of (icon_path_or_None, card_title, description) tuples, 2-4 items."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, intro)

    n = len(cards)
    gap = Inches(0.3)
    left0 = Inches(0.6)
    top = Inches(2.4)
    card_h = Inches(2.3)
    card_w = Emu(int((SLIDE_W - left0 * 2 - gap * (n - 1)) / n))
    left = left0
    for icon_path, card_title, desc in cards:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.adjustments[0] = 0.06
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_FILL
        card.line.color.rgb = NAVY
        card.line.width = Pt(1.25)
        card.shadow.inherit = False
        card.text_frame.clear()  # we'll place separate shapes instead of using the card's own text

        if icon_path and os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, left + Inches(0.3), top + Inches(0.3), height=Inches(0.9))

        tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.4), card_w - Inches(0.6), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        _run(p, card_title, 15, ACQUIA_BLUE, bold=True, font=FONT_HEADING)

        tb2 = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.95), card_w - Inches(0.6), Inches(1.3))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        _run(p2, desc, 12.5, BODY_GRAY)

        left = Emu(int(left + card_w + gap))

    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 3: Process / numbered step row, with arrows between steps
# ---------------------------------------------------------------------------
def add_process_steps_slide(prs, eyebrow, title, intro, steps, logo_variant="blue"):
    """steps: list of (number_str, icon_path_or_None, step_title, description) tuples."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, intro)

    n = len(steps)
    left0 = Inches(0.6)
    top = Inches(2.5)
    arrow_w = Inches(0.35)
    gap = Inches(0.15)
    col_w = Emu(int((SLIDE_W - left0 * 2 - arrow_w * (n - 1) - gap * (n - 1) * 2) / n))
    left = left0
    for i, (num, icon_path, step_title, desc) in enumerate(steps):
        tb = slide.shapes.add_textbox(left, top, col_w, Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        _run(p, num, 14, ACQUIA_BLUE, bold=True, font=FONT_HEADING)

        if icon_path and os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, left, top + Inches(0.45), height=Inches(0.9))

        tb2 = slide.shapes.add_textbox(left, top + Inches(1.5), col_w, Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]
        _run(p2, step_title, 14, NAVY, bold=True, font=FONT_HEADING)

        tb3 = slide.shapes.add_textbox(left, top + Inches(1.95), col_w, Inches(1.4))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        _run(p3, desc, 12.5, BODY_GRAY)

        left = Emu(int(left + col_w + gap))
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top + Inches(0.55), arrow_w, Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SKY_BLUE
            arrow.line.fill.background()
            left = Emu(int(left + arrow_w + gap))

    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 4: Flow / emphasis comparison (boxes connected by arrows, one highlighted)
# ---------------------------------------------------------------------------
def add_flow_emphasis_slide(prs, eyebrow, title, intro, boxes, highlight_index, logo_variant="blue"):
    """boxes: list of (label, description) tuples. highlight_index: which box is the solid-navy one."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, intro)

    n = len(boxes)
    left0 = Inches(0.6)
    top = Inches(2.6)
    box_h = Inches(2.6)
    arrow_w = Inches(0.4)
    gap = Inches(0.1)
    box_w = Emu(int((SLIDE_W - left0 * 2 - arrow_w * (n - 1) - gap * (n - 1) * 2) / n))
    left = left0
    for i, (label, desc) in enumerate(boxes):
        is_hl = (i == highlight_index)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.adjustments[0] = 0.06
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY if is_hl else WHITE
        box.line.color.rgb = NAVY
        box.line.width = Pt(1.25) if not is_hl else Pt(0)
        box.shadow.inherit = False
        box.text_frame.clear()

        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), box_w - Inches(0.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        _run(p, label, 15, WHITE if is_hl else NAVY, bold=True, font=FONT_HEADING)

        tb2 = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.9), box_w - Inches(0.5), Inches(1.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        _run(p2, desc, 12.5, RGBColor(0xE5, 0xEA, 0xF5) if is_hl else BODY_GRAY)

        left = Emu(int(left + box_w + gap))
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top + box_h / 2 - Inches(0.175), arrow_w, Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = SKY_BLUE
            arrow.line.fill.background()
            left = Emu(int(left + arrow_w + gap))

    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 5: Two-column highlight box (one navy-filled, one light card)
# ---------------------------------------------------------------------------
def add_two_column_highlight_slide(prs, eyebrow, title, intro,
                                    left_heading, left_bullets,
                                    right_heading, right_bullets, logo_variant="blue"):
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, intro)

    top = Inches(2.4)
    box_h = Inches(3.6)
    box_w = Inches(5.9)
    gap = Inches(0.5)
    left_x = Inches(0.6)
    right_x = Emu(int(left_x + box_w + gap))

    # Left: solid navy box
    lbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top, box_w, box_h)
    lbox.adjustments[0] = 0.05
    lbox.fill.solid(); lbox.fill.fore_color.rgb = NAVY
    lbox.line.fill.background(); lbox.shadow.inherit = False
    lbox.text_frame.clear()
    tb = slide.shapes.add_textbox(left_x + Inches(0.35), top + Inches(0.3), box_w - Inches(0.7), Inches(0.5))
    _run(tb.text_frame.paragraphs[0], left_heading, 16, SKY_BLUE, bold=True, font=FONT_HEADING)
    tb2 = slide.shapes.add_textbox(left_x + Inches(0.35), top + Inches(0.95), box_w - Inches(0.7), Inches(2.4))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for i, item in enumerate(left_bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        _run(p, f"-  {item}", 13, WHITE)

    # Right: light card box
    rbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top, box_w, box_h)
    rbox.adjustments[0] = 0.05
    rbox.fill.solid(); rbox.fill.fore_color.rgb = CARD_FILL
    rbox.line.color.rgb = NAVY; rbox.line.width = Pt(1.25); rbox.shadow.inherit = False
    rbox.text_frame.clear()
    tb3 = slide.shapes.add_textbox(right_x + Inches(0.35), top + Inches(0.3), box_w - Inches(0.7), Inches(0.5))
    _run(tb3.text_frame.paragraphs[0], right_heading, 16, ACQUIA_BLUE, bold=True, font=FONT_HEADING)
    tb4 = slide.shapes.add_textbox(right_x + Inches(0.35), top + Inches(0.95), box_w - Inches(0.7), Inches(2.4))
    tf4 = tb4.text_frame; tf4.word_wrap = True
    for i, item in enumerate(right_bullets):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.space_after = Pt(8)
        _run(p, f"-  {item}", 13, BODY_GRAY)

    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 6: Generic content slide (bullets + optional chart/image), for
# data findings, charts (build the chart PNG separately per the dataviz
# skill, then pass its path here), or screenshots.
# ---------------------------------------------------------------------------
def _render_bullet_list(tf, bullets, size, first=True):
    """bullets items are either a plain string (rendered as a "-  " bullet) or
    a (True, "Heading text") tuple (rendered bold, no dash, as a sub-heading
    line within the same textbox). Spacing scales down with font size so
    dense lists (8+ items) don't overflow their box - always check the
    rendered thumbnail on a dense slide rather than assuming it fit."""
    space_after = max(4, min(10, size * 0.6))
    space_before = max(2, min(6, size * 0.4))
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(space_after)
        if isinstance(item, tuple) and item[0] is True:
            p.space_before = Pt(space_before) if i > 0 else Pt(0)
            _run(p, item[1], size + 1, ACQUIA_BLUE, bold=True, font=FONT_HEADING)
        else:
            _run(p, f"-  {item}", size, BODY_GRAY)


def add_content_chart_slide(prs, eyebrow, title, intro, bullets=None, chart_path=None, footnote=None, logo_variant="blue", bullets_size=None):
    """bullets: list of strings, or (True, "Heading") tuples to insert a bold
    sub-heading line (see _render_bullet_list). Pass bullets_size explicitly
    (e.g. 11-12) for a dense list (8+ items incl. headings) sharing the
    narrower sidebar column - the default (13.5) will overflow past the
    footnote/logo on a long list. Always check the rendered thumbnail."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, intro)

    top = Inches(2.3)
    if chart_path and os.path.exists(chart_path):
        chart_w = Inches(7.6) if bullets else Inches(12.1)
        slide.shapes.add_picture(chart_path, Inches(0.6), top, width=chart_w)
        if bullets:
            tb = slide.shapes.add_textbox(Inches(8.4), top, Inches(4.3), Inches(4.55))
            tf = tb.text_frame; tf.word_wrap = True
            _render_bullet_list(tf, bullets, bullets_size or 13.5)
    elif bullets:
        tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(4.55))
        tf = tb.text_frame; tf.word_wrap = True
        _render_bullet_list(tf, bullets, bullets_size or 14)

    if footnote:
        add_footnote(slide, footnote)
    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 8: Stat card row (KPI headline numbers + signed deltas), for a
# quarterly/period metrics summary. Per dataviz's "stat tile" contract: a
# single headline number is not a chart.
# ---------------------------------------------------------------------------
def add_stat_cards_slide(prs, eyebrow, title, headline, cards, bullets=None, footnote=None, logo_variant="blue"):
    """cards: list of (value, delta_text_or_None, direction, label) tuples.
    direction: "good" (teal), "bad" (orange), or None (gray, e.g. "no data")."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, headline)

    n = len(cards)
    gap = Inches(0.25)
    left0 = Inches(0.6)
    top = Inches(2.0)
    card_h = Inches(1.7)
    card_w = Emu(int((SLIDE_W - left0 * 2 - gap * (n - 1)) / n))
    left = left0
    for value, delta_text, direction, label in cards:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.adjustments[0] = 0.08
        card.fill.solid(); card.fill.fore_color.rgb = CARD_FILL
        card.line.fill.background(); card.shadow.inherit = False
        card.text_frame.clear()

        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), card_w - Inches(0.4), Inches(0.55))
        _run(tb.text_frame.paragraphs[0], value, 26, NAVY, bold=True, font=FONT_HEADING)

        if delta_text:
            delta_color = TEAL if direction == "good" else (ORANGE if direction == "bad" else BODY_GRAY)
            tb2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.72), card_w - Inches(0.4), Inches(0.35))
            _run(tb2.text_frame.paragraphs[0], delta_text, 12.5, delta_color, bold=True)

        tb3 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.1), card_w - Inches(0.4), Inches(0.55))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        _run(tf3.paragraphs[0], label, 11.5, BODY_GRAY)

        left = Emu(int(left + card_w + gap))

    if bullets:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(4.05), Inches(12.1), Inches(2.7))
        tf = tb.text_frame; tf.word_wrap = True
        _render_bullet_list(tf, bullets, 14)

    if footnote:
        add_footnote(slide, footnote)
    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide


# ---------------------------------------------------------------------------
# Template 9: Numbered priority list (stacked cards with a left accent bar),
# for a prioritized action plan / roadmap closing slide.
# ---------------------------------------------------------------------------
def add_priority_list_slide(prs, eyebrow, title, headline, items, footnote=None, logo_variant="blue"):
    """items: list of (number, title, description, accent) tuples.
    accent: ACQUIA_BLUE (do now) or ORANGE (schedule later) etc."""
    slide = _blank(prs)
    _set_bg(slide, WHITE)
    _eyebrow_title_intro(slide, eyebrow, title, headline)

    top = Inches(2.05)
    row_h = Inches(0.86)
    gap = Inches(0.08)
    left = Inches(0.6)
    width = Inches(12.2)
    bar_w = Pt(4)

    for number, item_title, desc, accent in items:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, bar_w, row_h)
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background(); bar.shadow.inherit = False

        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(left + bar_w)), top, Emu(int(width - bar_w)), row_h)
        card.fill.solid(); card.fill.fore_color.rgb = CARD_FILL
        card.line.fill.background(); card.shadow.inherit = False
        card.text_frame.clear()

        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.06), width - Inches(0.5), Inches(0.32))
        _run(tb.text_frame.paragraphs[0], f"{number} · {item_title}", 13.5, NAVY, bold=True, font=FONT_HEADING)

        tb2 = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.38), width - Inches(0.5), Inches(0.45))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        _run(tf2.paragraphs[0], desc, 10.5, BODY_GRAY)

        top = Emu(int(top + row_h + gap))

    if footnote:
        add_footnote(slide, footnote)
    add_logo(slide, variant=logo_variant, top=Inches(6.85), left=Inches(11.3), width=Inches(1.3))
    return slide
